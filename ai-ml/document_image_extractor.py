# document_image_extractor.py

import os
import fitz  # PyMuPDF
import uuid
import logging
import io
from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# --------------------------------------------------
# UTILS
# --------------------------------------------------
def _ensure_dir(path):
    os.makedirs(path, exist_ok=True)


# --------------------------------------------------
# PDF IMAGE EXTRACTION
# --------------------------------------------------
def extract_images_from_pdf(file_path, output_dir="extracted_images/pdf"):
    extracted = []
    _ensure_dir(output_dir)

    try:
        with fitz.open(file_path) as pdf:
            for page_index, page in enumerate(pdf):
                images = page.get_images(full=True)

                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = pdf.extract_image(xref)

                    image_bytes = base_image.get("image")
                    ext = base_image.get("ext", "png")

                    if not image_bytes:
                        continue

                    filename = f"pdf_p{page_index+1}_{img_index+1}.{ext}"
                    img_path = os.path.join(output_dir, filename)

                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    extracted.append({
                        "source": "pdf",
                        "page": page_index + 1,
                        "path": img_path
                    })

        logger.info(
            f"PDF embedded images extracted: {len(extracted)} "
            f"from {os.path.basename(file_path)}"
        )
        return extracted

    except Exception as e:
        logger.error(
            f"PDF image extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# DOCX IMAGE EXTRACTION
# --------------------------------------------------
def extract_images_from_docx(file_path, output_dir="extracted_images/docx"):
    extracted = []
    _ensure_dir(output_dir)

    try:
        doc = Document(file_path)

        for rel in doc.part.rels.values():
            if "image" not in rel.target_ref:
                continue

            image_data = rel.target_part.blob
            img = Image.open(io.BytesIO(image_data)).convert("RGB")

            filename = f"docx_{uuid.uuid4().hex}.png"
            path = os.path.join(output_dir, filename)
            img.save(path)

            extracted.append({
                "source": "docx",
                "path": path
            })

        logger.info(
            f"DOCX embedded images extracted: {len(extracted)} "
            f"from {os.path.basename(file_path)}"
        )
        return extracted

    except Exception as e:
        logger.error(
            f"DOCX image extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# PPTX IMAGE EXTRACTION
# --------------------------------------------------
def extract_images_from_pptx(file_path, output_dir="extracted_images/pptx"):
    extracted = []
    _ensure_dir(output_dir)

    try:
        prs = Presentation(file_path)

        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                img_bytes = shape.image.blob
                img = Image.open(io.BytesIO(img_bytes)).convert("RGB")

                filename = f"pptx_s{slide_index+1}_{uuid.uuid4().hex}.png"
                path = os.path.join(output_dir, filename)
                img.save(path)

                extracted.append({
                    "source": "pptx",
                    "slide": slide_index + 1,
                    "path": path
                })

        logger.info(
            f"PPTX embedded images extracted: {len(extracted)} "
            f"from {os.path.basename(file_path)}"
        )
        return extracted

    except Exception as e:
        logger.error(
            f"PPTX image extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []
