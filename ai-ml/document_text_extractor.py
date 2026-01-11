# document_text_extractor.py

import os
import logging
import pdfplumber
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

# --------------------------------------------------
# PDF TEXT EXTRACTION
# --------------------------------------------------
def extract_text_from_pdf(file_path):
    pages = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                pages.append({
                    "page": i + 1,
                    "text": text.strip()
                })

        logger.info(f"PDF text extracted: {len(pages)} pages")
        return pages

    except Exception as e:
        logger.error(f"PDF extraction failed: {e}", exc_info=True)
        return []


# --------------------------------------------------
# DOCX TEXT EXTRACTION
# --------------------------------------------------
def extract_text_from_docx(file_path):
    blocks = []

    try:
        doc = Document(file_path)
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                blocks.append({
                    "paragraph": i + 1,
                    "text": text
                })

        logger.info(f"DOCX text extracted: {len(blocks)} paragraphs")
        return blocks

    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}", exc_info=True)
        return []


# --------------------------------------------------
# PPTX TEXT EXTRACTION
# --------------------------------------------------
def extract_text_from_pptx(file_path):
    slides = []

    try:
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        slide_text.append(txt)

            slides.append({
                "slide": i + 1,
                "text": "\n".join(slide_text)
            })

        logger.info(f"PPTX text extracted: {len(slides)} slides")
        return slides

    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}", exc_info=True)
        return []


# --------------------------------------------------
# TXT TEXT EXTRACTION
# --------------------------------------------------
def extract_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        logger.info("TXT text extracted")
        return [{
            "text": text.strip()
        }]

    except Exception as e:
        logger.error(f"TXT extraction failed: {e}", exc_info=True)
        return []


# --------------------------------------------------
# ROUTER
# --------------------------------------------------
def extract_document_text(file_path, doc_type):
    """
    Unified router for document text extraction
    """

    if doc_type == "pdf":
        return extract_text_from_pdf(file_path)

    if doc_type == "docx":
        return extract_text_from_docx(file_path)

    if doc_type == "pptx":
        return extract_text_from_pptx(file_path)

    if doc_type == "txt":
        return extract_text_from_txt(file_path)

    logger.warning(f"No extractor for document type: {doc_type}")
    return []
