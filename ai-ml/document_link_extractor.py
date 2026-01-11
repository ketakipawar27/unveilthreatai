# document_link_extractor.py

import re
import os
import logging
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

# --------------------------------------------------
# REGEX PATTERNS (NON-CAPTURING, SAFE)
# --------------------------------------------------
URL_PATTERN = re.compile(
    r"\bhttps?://[^\s<>\"]+|\bwww\.[^\s<>\"]+",
    re.IGNORECASE
)

EMAIL_PATTERN = re.compile(
    r"\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+\b"
)

IP_URL_PATTERN = re.compile(
    r"\bhttps?://(?:\d{1,3}\.){3}\d{1,3}(?::\d+)?\b"
)

# --------------------------------------------------
# CORE UTILITY
# --------------------------------------------------
def _normalize_link(link: str) -> str:
    """Ensure consistent URL format"""
    link = link.strip()
    if link.startswith("www."):
        return "http://" + link
    return link


def _extract_links_from_text(text):
    if not text:
        return []

    links = set()

    for match in URL_PATTERN.findall(text):
        links.add(_normalize_link(match))

    for match in EMAIL_PATTERN.findall(text):
        links.add(f"mailto:{match}")

    for match in IP_URL_PATTERN.findall(text):
        links.add(match)

    return list(links)


# --------------------------------------------------
# PDF LINK EXTRACTION
# --------------------------------------------------
def extract_links_from_pdf(file_path):
    links = []
    seen = set()

    try:
        with fitz.open(file_path) as pdf:
            for page_num, page in enumerate(pdf):
                text = page.get_text() or ""
                found = _extract_links_from_text(text)

                for link in found:
                    if link in seen:
                        continue
                    seen.add(link)

                    links.append({
                        "type": "link",
                        "source": "pdf",
                        "page": page_num + 1,
                        "value": link
                    })

        logger.info(
            f"PDF links extracted: {len(links)} "
            f"from {os.path.basename(file_path)}"
        )
        return links

    except Exception as e:
        logger.error(
            f"PDF link extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# DOCX LINK EXTRACTION
# --------------------------------------------------
def extract_links_from_docx(file_path):
    links = []
    seen = set()

    try:
        doc = Document(file_path)

        for para_index, para in enumerate(doc.paragraphs):
            found = _extract_links_from_text(para.text)

            for link in found:
                if link in seen:
                    continue
                seen.add(link)

                links.append({
                    "type": "link",
                    "source": "docx",
                    "paragraph": para_index + 1,
                    "value": link
                })

        logger.info(
            f"DOCX links extracted: {len(links)} "
            f"from {os.path.basename(file_path)}"
        )
        return links

    except Exception as e:
        logger.error(
            f"DOCX link extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# PPTX LINK EXTRACTION
# --------------------------------------------------
def extract_links_from_pptx(file_path):
    links = []
    seen = set()

    try:
        prs = Presentation(file_path)

        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                found = _extract_links_from_text(shape.text)

                for link in found:
                    if link in seen:
                        continue
                    seen.add(link)

                    links.append({
                        "type": "link",
                        "source": "pptx",
                        "slide": slide_index + 1,
                        "value": link
                    })

        logger.info(
            f"PPTX links extracted: {len(links)} "
            f"from {os.path.basename(file_path)}"
        )
        return links

    except Exception as e:
        logger.error(
            f"PPTX link extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# GENERIC TEXT FILE
# --------------------------------------------------
def extract_links_from_text_file(file_path):
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        found = _extract_links_from_text(text)
        seen = set()

        return [
            {
                "type": "link",
                "source": "text",
                "value": link
            }
            for link in found
            if not (link in seen or seen.add(link))
        ]

    except Exception as e:
        logger.error(
            f"Text link extraction failed for {file_path}: {e}",
            exc_info=True
        )
        return []


# --------------------------------------------------
# MASTER ROUTER
# --------------------------------------------------
def extract_links(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_links_from_pdf(file_path)
    elif ext == ".docx":
        return extract_links_from_docx(file_path)
    elif ext == ".pptx":
        return extract_links_from_pptx(file_path)
    elif ext in [".txt", ".md"]:
        return extract_links_from_text_file(file_path)
    else:
        logger.warning(f"Unsupported file for link extraction: {ext}")
        return []
